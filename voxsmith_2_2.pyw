import os
import sys
import json
import re
import keyring
import threading
import time
import hashlib
import subprocess
import tempfile
import traceback
import shutil
import tkinter as tk
from tkinter import filedialog, messagebox, simpledialog, ttk
import customtkinter as ctk
import requests
import certifi
from urllib3.util.retry import Retry
from requests.adapters import HTTPAdapter
from pptx import Presentation

# --- Security: redaction helpers (Step 5) ---
REDACTION_PATTERNS = [
    (re.compile(r"(xi-api-key:\s*)([A-Za-z0-9_\-]{10,})", re.IGNORECASE), r"\1[REDACTED]"),
    (re.compile(r"(Authorization:\s*Bearer\s+)([_A-Za-z0-9\.\-]{10,})", re.IGNORECASE), r"\1[REDACTED]"),
    (re.compile(r"([?&](?:xi-api-key|api_key|apikey|token)=)([^&\s]{6,})", re.IGNORECASE), r"\1[REDACTED]"),
    (re.compile(r"[A-Fa-f0-9]{8}-[A-Fa-f0-9]{4}-[A-Fa-f0-9]{4}-[A-Fa-f0-9]{4}-[A-Fa-f0-9]{12}"), "[REDACTED]"),
    (re.compile(r"[A-Za-z0-9_\-]{32,}"), "[REDACTED]"),
    (re.compile(r"[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}"), "[REDACTED]"),
]
import logging
from logging.handlers import RotatingFileHandler
import atexit
from datetime import datetime

# --- Phase C: Network Telemetry Helper ---
import urllib.parse

def _voxsmith_http(method: str, url: str, **kwargs):
    """Wrap outbound HTTP with timing and lightweight header logging.
    Returns (response). Raises requests exceptions like the underlying call.
    """
    logger = logging.getLogger("voxsmith")
    t0 = time.perf_counter()
    try:
        sess = make_voxsmith_session()
    except Exception:
        sess = requests.Session()
    if "timeout" not in kwargs:
        kwargs["timeout"] = NET_TIMEOUT if "NET_TIMEOUT" in globals() else 120
    resp = None
    err = None
    try:
        headers = kwargs.get("headers") or {}
        headers.setdefault("User-Agent", USER_AGENT)
        kwargs["headers"] = headers
        resp = sess.request(method.upper(), url, **kwargs)
        return resp
    except Exception as e:
        err = e
        raise
    finally:
        t1 = time.perf_counter()
        elapsed_ms = int((t1 - t0) * 1000)
        try:
            parts = urllib.parse.urlsplit(url)
            path_only = parts.path
            rid = None
            try:
                rid = (resp.headers.get("x-request-id") if resp is not None else None) or \
                      (resp.headers.get("x-amzn-requestid") if resp is not None else None) or \
                      (resp.headers.get("request-id") if resp is not None else None)
            except Exception:
                rid = None
            date_hdr = None
            try:
                date_hdr = resp.headers.get("date") if resp is not None else None
            except Exception:
                date_hdr = None
            retry_after = None
            try:
                retry_after = resp.headers.get("retry-after") if resp is not None else None
            except Exception:
                retry_after = None
            size_bytes = None
            try:
                if resp is not None:
                    size_bytes = int(resp.headers.get("content-length") or 0) or (len(resp.content) if getattr(resp, "content", None) is not None else 0)
            except Exception:
                size_bytes = None
            status = getattr(resp, "status_code", None)
            try:
                _handle_auth_log(resp)
            except Exception:
                pass
            log_msg = {
                "evt": "net",
                "method": method.upper(),
                "path": path_only,
                "status": status,
                "ms": elapsed_ms,
                "bytes": size_bytes,
                "date": date_hdr,
                "rid": rid,
                "retry_after": retry_after,
                "ua_ver": globals().get("APP_VERSION", "v?")
            }
            try:
                logger.info(_redact(json.dumps(log_msg, ensure_ascii=False))) if NET_VERBOSE else logger.info(_redact(f"NET {method.upper()} {path_only} status={status} ms={elapsed_ms} bytes={size_bytes}"))
            except Exception:
                logger.info(_redact(f"NET {method.upper()} {path_only} status={status} ms={elapsed_ms} bytes={size_bytes} rid={rid} date={date_hdr} retry_after={retry_after}"))
        except Exception:
            pass


# [Phase C] network guardrails
from voxsecurity.allowlist import make_voxsmith_session, DomainNotAllowed
from voxsecurity.checksum_verify import verify_self
import time

from voxsecurity.checksum_verify import verify_self

# Single outbound session restricted to approved domains
VOX_SESSION = make_voxsmith_session()  # api.elevenlabs.io, update.voxsmith.app
# Back-compat: reuse existing session variable
_VOX_SESSION = VOX_SESSION

# Phase C: allowlisted HTTP session
VOX_SESSION = make_voxsmith_session()  # restricts outbound domains


def _redact(s: str) -> str:
    try:
        out = str(s)
        for pat, repl in REDACTION_PATTERNS:
            out = re.sub(pat, repl, out)
        return out
    except Exception:
        return str(s)

class RedactingFilter(logging.Filter):
    def filter(self, record):
        try:
            msg = record.getMessage()
            record.msg = _redact(msg)
            record.args = None
        except Exception:
            pass
        return True

# --- App paths & file logging (Phase A) ---
APP_NAME = "Voxsmith"
def _get_app_paths():
    base = os.getenv("LOCALAPPDATA") or os.getenv("APPDATA") or tempfile.gettempdir()
    root = os.path.join(base, APP_NAME)
    logs = os.path.join(root, "logs")
    tempd = os.path.join(root, "temp")
    for d in (root, logs, tempd):
        try:
            os.makedirs(d, exist_ok=True)
        except Exception:
            pass
    return {"root": root, "logs": logs, "temp": tempd}

def _logout_flag_path():
    p = os.path.join(_get_app_paths()["root"], "logout.flag")
    return p

def _set_logout_flag():
    try:
        with open(_logout_flag_path(), "w", encoding="utf-8") as f:
            f.write("logout")
    except Exception:
        pass

def _clear_logout_flag():
    try:
        p = _logout_flag_path()
        if os.path.exists(p):
            os.remove(p)
    except Exception:
        pass

def _is_logged_out() -> bool:
    try:
        return os.path.exists(_logout_flag_path())
    except Exception:
        return False


def _configure_file_logger(logger_name="voxsmith"):
    logger = logging.getLogger(logger_name)
    # avoid duplicate rotating handlers
    if any(isinstance(h, RotatingFileHandler) for h in logger.handlers):
        return logger
    try:
        paths = _get_app_paths()
        log_path = os.path.join(paths["logs"], "voxsmith.log")
        fh = RotatingFileHandler(log_path, maxBytes=512*1024, backupCount=3, encoding="utf-8")
        fmt = logging.Formatter("[%(asctime)s] %(levelname)s %(name)s: %(message)s")
        fh.setFormatter(fmt)
        try:
            fh.addFilter(RedactingFilter())
        except Exception:
            pass
        logger.addHandler(fh)
    except Exception:
        pass
    return logger



import voxattach
import voxanimate

def attach_audio_for_slide(deck_path: str, slide_index_1based: int, src_audio: str, out_audio: str):
    try:
        logger = logging.getLogger("voxsmith")
        # Prefer explicit attach_to_slide if available; fall back to attach_or_skip
        if hasattr(voxattach, "attach_to_slide"):
            result = voxattach.attach_to_slide(deck_path, slide_index_1based, out_audio)
        else:
            result = voxattach.attach_or_skip(deck_path, slide_index_1based, src_audio, out_audio)
        try:
            logger.info(_redact(f"ATTACH slide={slide_index_1based} deck={os.path.basename(deck_path)} result={result}"))
        except Exception:
            pass
        print(f"[voxsmith] attach slide {slide_index_1based}: {result}")
        return result
    except Exception as e:
        traceback.print_exc()
        try:
            logging.getLogger("voxsmith").error(_redact(f"ATTACH exception slide={slide_index_1based}: {e}"))
        except Exception:
            pass
        print(f"[voxsmith] attach failed on slide {slide_index_1based}: {e}")
        return {"processed": False, "attached": False, "reason": "exception", "error": str(e)}

# winsound is Windows-only; keep optional
try:
    import winsound
except Exception:
    winsound = None

APP_NAME = "Voxsmith 2"
APP_VERSION = "v2.2"
# --- Phase C: UA/version hygiene & verbosity toggle ---
NET_VERBOSE = False
def _compute_build_hash() -> str:
    try:
        import hashlib
        p = __file__
        h = hashlib.sha1()
        with open(p, "rb") as f:
            for chunk in iter(lambda: f.read(8192), b""):
                h.update(chunk)
        return h.hexdigest()[:7]
    except Exception:
        return "local"
BUILD_HASH = _compute_build_hash()
USER_AGENT = f"Voxsmith/{APP_VERSION}+{BUILD_HASH}"
# --- Network hygiene (Step 2) ---
NET_MAX_ATTEMPTS = 3  # 1 initial + 2 retries
NET_BACKOFF_BASE = 0.75  # seconds; exponential backoff
NET_TIMEOUT = 120  # seconds per request
try:
    _VOX_SESSION = VOX_SESSION

    try:
        # Configure retries
        retry_kwargs = dict(total=3, backoff_factor=0.5, status_forcelist=(429,500,502,503,504), raise_on_status=False)
        try:
            retry = Retry(allowed_methods=frozenset(["GET","PUT","DELETE","HEAD","OPTIONS"]), **retry_kwargs)
        except TypeError:
            retry = Retry(method_whitelist=frozenset(["GET","PUT","DELETE","HEAD","OPTIONS"]), **retry_kwargs)
        adapter = HTTPAdapter(max_retries=retry, pool_connections=10, pool_maxsize=10)
        _VOX_SESSION.mount("https://", adapter)
        _VOX_SESSION.mount("http://", adapter)
        # Pin CA bundle for frozen builds
        _VOX_SESSION.verify = certifi.where()
        # Set a default UA
        _VOX_SESSION.headers.update({"User-Agent": "Voxsmith/2.14.2 (+https://donburnside.com/voxsmith-2/)"})
    except Exception:
        pass
except Exception:
    _VOX_SESSION = None

TARGET_SAMPLE_RATE = "44100"
TARGET_CODEC = "pcm_s16le"
TARGET_CHANNELS = "2"

DEFAULT_API_KEY = ""
DEFAULT_INPUT_FILE = ""
DEFAULT_OUTPUT_DIR = ""
DEFAULT_PREVIEW_TEXT = "This is a quick voice preview."
DEFAULT_SLIDE_RANGE = ""
DEFAULT_FIXED_ONLY = False
DEFAULT_MAKE_COPY = True
DEFAULT_HIDE_ICON = True

def _ensure_local_ffmpeg_on_path():
    try:
        if getattr(sys, "frozen", False):
            exe_dir = os.path.dirname(sys.executable)
        else:
            exe_dir = os.path.dirname(os.path.abspath(__file__))
        local_ffmpeg_dir = os.path.join(exe_dir, "ffmpeg")
        if os.path.isdir(local_ffmpeg_dir):
            os.environ["PATH"] = local_ffmpeg_dir + os.pathsep + os.environ.get("PATH", "")
    except Exception:
        pass

def normalize_audio(input_file: str, output_file: str):
    cmd = [
        "ffmpeg", "-y",
        "-i", input_file,
        "-ar", TARGET_SAMPLE_RATE,
        "-ac", TARGET_CHANNELS,
        "-acodec", TARGET_CODEC,
        output_file
    ]
    subprocess.run(cmd, check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    return output_file

_ensure_local_ffmpeg_on_path()

def get_settings_dir() -> str:
    if os.name == "nt":
        base = os.getenv("APPDATA") or os.path.expanduser("~")
        return os.path.join(base, APP_NAME)
    elif sys.platform == "darwin":
        return os.path.expanduser(f"~/Library/Application Support/{APP_NAME}")
    else:
        return os.path.expanduser(f"~/.config/{APP_NAME}")

SETTINGS_DIR = get_settings_dir()
LOGS_DIR = os.path.join(SETTINGS_DIR, 'logs')
SETTINGS_FILE = os.path.join(SETTINGS_DIR, "settings.json")

def safe_ensure_dir(path: str) -> bool:
    try:
        os.makedirs(path, exist_ok=True)
        return True
    except Exception:
        return False

def load_settings():
    try:
        if os.path.exists(SETTINGS_FILE):
            with open(SETTINGS_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
    except Exception:
        pass
    return {}

def save_settings(**kwargs) -> bool:
    data = load_settings()
    data.update(kwargs)
    if safe_ensure_dir(SETTINGS_DIR):
        try:
            with open(SETTINGS_FILE, "w", encoding="utf-8") as f:
                json.dump(data, f, indent=2)
            return True
        except Exception:
            pass
    fallback_dir = os.path.join(tempfile.gettempdir(), APP_NAME)
    if safe_ensure_dir(fallback_dir):
        try:
            with open(os.path.join(fallback_dir, "settings.json"), "w", encoding="utf-8") as f:
                json.dump(data, f, indent=2)
            return True
        except Exception:
            pass
    return False

# === Secure API key helpers (Windows Credential Manager via keyring) ===
def get_api_key() -> str:
    # Prefer environment variable if set for CI/testing, else OS vault
    env = os.getenv("ELEVENLABS_API_KEY", "").strip()
    if env:
        return env
    try:
        val = keyring.get_password(APP_NAME, "elevenlabs") or ""
        return val.strip()
    except Exception:
        return ""

def set_api_key(k: str) -> None:
    try:
        keyring.set_password(APP_NAME, "elevenlabs", (k or "").strip())
    except Exception:
        pass

def delete_api_key() -> None:
    try:
        keyring.delete_password(APP_NAME, "elevenlabs")
    except Exception:
        # If it doesn't exist or backend errors, ignore
        pass

# === Voice caching helpers ===
def get_voices_cache_path() -> str:
    """Get path to voices cache file."""
    try:
        settings_dir = get_settings_dir()
        return os.path.join(settings_dir, "voices_cache.json")
    except Exception:
        return ""

def load_voice_cache() -> dict:
    """Load cached voices if valid (age < 1 week, key matches)."""
    try:
        path = get_voices_cache_path()
        if not path or not os.path.exists(path):
            return None
            
        with open(path, 'r', encoding='utf-8') as f:
            cache = json.load(f)
        
        # Verify cache has required fields
        if not all(k in cache for k in ['api_key_hash', 'last_updated', 'voices']):
            return None
        
        # Verify API key matches
        current_key = get_api_key()
        if not current_key:
            return None
        key_hash = hashlib.sha256(current_key.encode('utf-8')).hexdigest()
        
        if cache.get('api_key_hash') != key_hash:
            return None  # Different key, invalidate cache
            
        # Check age (max 1 week = 604800 seconds)
        try:
            updated = datetime.fromisoformat(cache['last_updated'])
            age_seconds = (datetime.now() - updated).total_seconds()
            
            if age_seconds > 604800:  # 1 week
                return None  # Too old
        except Exception:
            return None
            
        return cache
    except Exception:
        return None

def save_voice_cache(voices: list) -> None:
    """Save voices list to cache with timestamp and key hash."""
    try:
        current_key = get_api_key()
        if not current_key:
            return
        key_hash = hashlib.sha256(current_key.encode('utf-8')).hexdigest()
        
        cache = {
            "version": "1.0",
            "api_key_hash": key_hash,
            "last_updated": datetime.now().isoformat(),
            "voices": voices
        }
        
        path = get_voices_cache_path()
        if not path:
            return
            
        safe_ensure_dir(get_settings_dir())
        with open(path, 'w', encoding='utf-8') as f:
            json.dump(cache, f, indent=2)
    except Exception:
        pass

def delete_voice_cache() -> None:
    """Delete voices cache (called on logout)."""
    try:
        path = get_voices_cache_path()
        if path and os.path.exists(path):
            os.remove(path)
    except Exception:
        pass

def voices_changed(cached_voices: list, fresh_voices: list) -> bool:
    """Check if voice list actually changed."""
    try:
        if not cached_voices:
            return True  # No cache, definitely changed
        
        # Build dictionaries for comparison
        cached_dict = {v['voice_id']: v['name'] for v in cached_voices if isinstance(v, dict) and 'voice_id' in v and 'name' in v}
        fresh_dict = {vid: nm for nm, vid in fresh_voices}  # fresh_voices is list of (name, voice_id) tuples
        
        # Check for differences
        if set(cached_dict.keys()) != set(fresh_dict.keys()):
            return True  # Voice IDs changed
            
        if cached_dict != fresh_dict:
            return True  # Names changed
            
        return False  # Identical
    except Exception:
        return True  # On error, assume changed

def safe_path(p: str) -> str:
    try:
        return os.path.basename(p) if isinstance(p, str) else p
    except Exception:
        return p

# --- Phase C: Checksum helpers ---
def _sha256_bytes(b: bytes) -> str:
    try:
        h = hashlib.sha256(); h.update(b); return h.hexdigest()
    except Exception:
        return ""
def _sha256_file(p: str) -> str:
    try:
        with open(p, 'rb') as f:
            h = hashlib.sha256()
            for chunk in iter(lambda: f.read(8192), b''):
                h.update(chunk)
            return h.hexdigest()
    except Exception:
        return ""



def log_line(w, msg):
    # Verbosity-aware logging with minimal noise by default
    try:
        verbose = False
        v = getattr(w, '_verbose_var', None)
        if v is not None:
            try:
                verbose = bool(v.get())
            except Exception:
                verbose = False
        if not verbose:
            allowed_prefixes = ('X','*','>','OK','i','API error','Network error')
            if not any(str(msg).startswith(p) for p in allowed_prefixes):
                return
            # Light path redaction when not verbose
            try:
                msg = re.sub(r'([A-Za-z]:\\[^\s]+|/[^\s]+)', lambda m: os.path.basename(m.group(0)), str(msg))
            except Exception:
                pass
    except Exception:
        pass
    # Security redaction for secrets/tokens
    try:
        msg = _redact(msg)
    except Exception:
        pass
    try:
        w.configure(state="normal"); w.insert("end", str(msg) + "\n"); w.see("end"); w.configure(state="disabled")
    except Exception:
        pass

def run_hidden(cmd_list):
    kwargs = dict(stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    if os.name == "nt":
        try:
            kwargs["creationflags"] = subprocess.CREATE_NO_WINDOW
        except AttributeError:
            si = subprocess.STARTUPINFO(); si.dwFlags = subprocess.STARTF_USESHOWWINDOW
            kwargs["startupinfo"] = si
    return subprocess.run(cmd_list, **kwargs)

def run_ffmpeg_quiet(cmd_list):
    return run_hidden(cmd_list + ["-hide_banner", "-loglevel", "error"])

def test_ffmpeg_available():
    try:
        cp = run_hidden(["ffmpeg", "-version"])
        return cp.returncode == 0
    except FileNotFoundError:
        return False

def open_folder(path: str):
    try:
        if os.name == "nt": os.startfile(path)
        elif sys.platform == "darwin": run_hidden(["open", path])
        else: run_hidden(["xdg-open", path])
    except Exception:
        pass

def pretty_api_error(resp):
    try:
        data = resp.json()
    except Exception:
        data = None
    msg = None
    if isinstance(data, dict):
        d = data.get("detail")
        if isinstance(d, dict):
            msg = d.get("message") or d.get("error")
        elif isinstance(d, list) and d:
            msg = d[0].get("message") if isinstance(d[0], dict) else str(d[0])
        elif isinstance(d, str):
            msg = d
        if not msg:
            msg = data.get("message") or data.get("error")
    if not msg:
        msg = resp.text.strip()
    return f"HTTP {resp.status_code}: {msg}" if msg else f"HTTP {resp.status_code}"

def fetch_voices(api_key: str):
    if not api_key.strip():
        raise ValueError("Missing API key")
    url = "https://api.elevenlabs.io/v1/voices"
    h = {"xi-api-key": api_key}
    attempts = 0
    last_err = None
    while attempts < NET_MAX_ATTEMPTS:
        attempts += 1
        try:
            sess = VOX_SESSION
            resp = sess.get(url, headers=h, timeout=NET_TIMEOUT)
            if resp.status_code == 200:
                data = resp.json()
                out = []
                for v in data.get("voices", []):
                    nm = (v.get("name") or "").strip() or "(unnamed voice)"
                    vid = v.get("voice_id") or ""
                    if vid:
                        out.append((nm, vid))
                out.sort(key=lambda x: x[0].lower())
                return out
            # Retry on transient server errors
            if 500 <= resp.status_code < 600 and attempts < NET_MAX_ATTEMPTS:
                time.sleep(NET_BACKOFF_BASE * (2 ** (attempts-1)))
                continue
            # Non-retryable API error
            raise RuntimeError(pretty_api_error(resp))
        except requests.RequestException as e:
            last_err = e
            if attempts < NET_MAX_ATTEMPTS:
                time.sleep(NET_BACKOFF_BASE * (2 ** (attempts-1)))
                continue
            raise RuntimeError(f"Network error: {e}") from e

def select_slides(total: int, spec: str):
    if not spec.strip():
        return list(range(1, total + 1))
    s = set()
    for part in [p.strip() for p in spec.split(',') if p.strip()]:
        if '-' in part:
            a, b = part.split('-', 1)
            try:
                st = int(a) if a else 1
            except Exception:
                st = 1
            try:
                en = int(b) if b else total
            except Exception:
                en = total
            st = max(1, st); en = min(total, en)
            if st <= en:
                s.update(range(st, en + 1))
        else:
            try:
                n = int(part)
                if 1 <= n <= total:
                    s.add(n)
            except Exception:
                pass
    return sorted(s)

class PreviewPlayer:
    def __init__(self, log_widget, preview_btn, stop_btn, get_preview_text):
        self.log_widget = log_widget
        self.preview_btn = preview_btn
        self.stop_btn = stop_btn  # Can be None if no stop button
        self.get_preview_text = get_preview_text
        self._thread = None
        self._temp_path = None
        self._lock = threading.Lock()

    def is_running(self):
        return self._thread is not None and self._thread.is_alive()

    def stop(self):
        with self._lock:
            if winsound and sys.platform.startswith("win"):
                try:
                    winsound.PlaySound(None, 0)
                except Exception:
                    pass
            self._cleanup_temp()

    def _cleanup_temp(self):
        with self._lock:
            if self._temp_path and os.path.exists(self._temp_path):
                try:
                    os.remove(self._temp_path)
                except Exception:
                    pass
            self._temp_path = None

    def _play_wav_sync(self, p):
        if sys.platform.startswith("win") and winsound:
            try:
                winsound.PlaySound(p, winsound.SND_FILENAME)
                return
            except Exception:
                pass
        for cmd in (["afplay", p], ["aplay", p], ["paplay", p], ["ffplay", "-autoexit", "-nodisp", "-loglevel", "error", p]):
            try:
                run_hidden(cmd)
                return
            except Exception:
                continue

    def preview(self, api_key: str, voice_id: str):
        if self.is_running():
            return

        def worker():
            self.preview_btn.configure(state="disabled")
            if self.stop_btn:
                self.stop_btn.configure(state="normal")
            try:
                t = self.get_preview_text().strip() or "This is a quick voice preview."
                if not api_key.strip():
                    log_line(self.log_widget, "X Missing API Key for preview.")
                    messagebox.showerror("Preview", "Enter API Key.")
                    return
                if not voice_id.strip():
                    log_line(self.log_widget, "X Missing Voice selection for preview.")
                    messagebox.showerror("Preview", "Choose a voice.")
                    return
                url = f"https://api.elevenlabs.io/v1/text-to-speech/{voice_id}"
                h = {"xi-api-key": api_key, "Content-Type": "application/json", "Accept": "audio/wav"}
                payload = {"text": t, "output_format": "wav", "voice_settings": {"stability": 0.5, "similarity_boost": 0.7}}
                log_line(self.log_widget, "i Requesting preview...")
                sess = VOX_SESSION
                resp = sess.post(url, headers=h, json=payload, timeout=30)
                if resp.status_code != 200:
                    msg = pretty_api_error(resp)
                    log_line(self.log_widget, f"X Preview failed: {msg}")
                    messagebox.showerror("Preview Error", msg)
                    return
                ct = (resp.headers.get("Content-Type") or "").lower()
                ext = ".mp3" if ("mpeg" in ct or "mp3" in ct) else ".wav"
                with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as src:
                    src.write(resp.content); src_path = src.name
                with tempfile.NamedTemporaryFile(delete=False, suffix=".wav") as dst:
                    dst_path = dst.name
                run_ffmpeg_quiet(["ffmpeg","-y","-i",src_path,"-acodec","pcm_s16le","-ar","44100","-ac","2",dst_path])
                try:
                    os.remove(src_path)
                except Exception:
                    pass
                self._temp_path = dst_path
                log_line(self.log_widget, "> Playing preview...")
                self._play_wav_sync(dst_path)
                log_line(self.log_widget, "OK Preview finished.")
            except Exception as e:
                log_line(self.log_widget, f"X Preview error: {e}")
            finally:
                self._cleanup_temp()
                self.preview_btn.configure(state="normal")
                if self.stop_btn:
                    self.stop_btn.configure(state="disabled")

        self._thread = threading.Thread(target=worker, daemon=True)
        self._thread.start()

_SINGLE_LOCK = None

def _check_single_instance():
    if os.name != "nt":
        return False, None
    import msvcrt
    lock_path = os.path.join(os.getenv("LOCALAPPDATA") or os.getcwd(), f"{APP_NAME}.lock")
    os.makedirs(os.path.dirname(lock_path), exist_ok=True)
    lock_file = open(lock_path, "a")
    try:
        msvcrt.locking(lock_file.fileno(), msvcrt.LK_NBLCK, 1)
        return False, lock_file
    except OSError:
        try:
            lock_file.close()
        except Exception:
            pass
        return True, None

def _release_single_instance(lock_file):
    if not lock_file or os.name != "nt":
        return
    import msvcrt
    try:
        lock_file.seek(0)
        msvcrt.locking(lock_file.fileno(), msvcrt.LK_UNLCK, 1)
    except Exception:
        pass
    try:
        lock_file.close()
    except Exception:
        pass

# --- Tooltip helper class ---
class ToolTip:
    """Simple tooltip for CustomTkinter widgets"""
    def __init__(self, widget, text):
        self.widget = widget
        self.text = text
        self.tooltip = None
        self.widget.bind("<Enter>", self.show_tooltip)
        self.widget.bind("<Leave>", self.hide_tooltip)
    
    def show_tooltip(self, event=None):
        if self.tooltip:
            return
        x = self.widget.winfo_rootx() + 10
        y = self.widget.winfo_rooty() + self.widget.winfo_height() + 5
        self.tooltip = tk.Toplevel(self.widget)
        self.tooltip.wm_overrideredirect(True)
        self.tooltip.wm_geometry(f"+{x}+{y}")
        label = tk.Label(self.tooltip, text=self.text, 
                        background="#ffffe0", foreground="#000000",
                        relief="solid", borderwidth=1,
                        font=("Open Sans", 10), padx=8, pady=4)
        label.pack()
    
    def hide_tooltip(self, event=None):
        if self.tooltip:
            self.tooltip.destroy()
            self.tooltip = None

# --- helper for remove detection ---
MSO_MEDIA = 16  # MsoShapeType.msoMedia

def _name_looks_like_ng_audio(name: str) -> bool:
    n = (name or '').lower()
    # match our common names: slide01, slide_01, slide001, 01, 001
    if n.startswith('slide'):
        return True
    if n and n[0].isdigit() and len(n) in (2, 3):
        return True
    return False

# ------------------------------------------------------------
# Slide text extraction for "### Read Slide" marker
# ------------------------------------------------------------

def extract_slide_text(slide):
    """
    Extract text from a slide's shapes and text boxes (excluding title).
    Returns text in natural reading order: top-to-bottom, left-to-right.
    
    Excludes:
    - Title shapes (first shape with title placeholder)
    - Grouped objects
    - Tables
    
    Returns:
        str: Extracted text with newlines between shapes
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
    
    text_items = []
    
    # Find the title placeholder to exclude it
    title_shape_id = None
    for shape in slide.shapes:
        if shape.is_placeholder:
            try:
                if shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                    title_shape_id = shape.shape_id
                    break
            except:
                pass
    
    # Collect text from shapes with their positions
    for shape in slide.shapes:
        # Skip title
        if shape.shape_id == title_shape_id:
            continue
            
        # Skip grouped objects (they're part of a group)
        try:
            if hasattr(shape, 'group_items'):
                # This is a group shape itself - skip individual processing
                continue
        except:
            pass
        
        # Skip tables
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            continue
        
        # Extract text from text frames (shapes and text boxes)
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                # Store with position for sorting
                top = shape.top
                left = shape.left
                text_items.append((top, left, text))
    
    # Sort by position: top-to-bottom (primary), then left-to-right (secondary)
    text_items.sort(key=lambda x: (x[0], x[1]))
    
    # Extract just the text in sorted order
    extracted_text = '\n'.join(item[2] for item in text_items)
    
    return extracted_text

# ------------------------------------------------------------
# PowerPoint attach routine that removes previous audio for each slide
# ------------------------------------------------------------

def generate_narration(api_key, voice_id, input_file, output_dir, fixed_only, slide_range_spec, cancel_event,
                       log_widget, start_button, cancel_button, audio_only=False):

    def worker():
        pp_app = None
        pp_pres = None
        try:
            # Prepare session-scoped manifest path in app logs directory
            try:
                safe_ensure_dir(LOGS_DIR)
            except Exception:
                pass
            deck_base = os.path.splitext(os.path.basename(input_file))[0]
            session_ts = time.strftime('%Y%m%d_%H%M%S')
            session_manifest_path = os.path.join(LOGS_DIR, f"{deck_base}_{session_ts}_manifest.json")
            if not api_key.strip():
                log_line(log_widget, "X Missing API Key."); messagebox.showerror("Error","Enter API Key."); return
            if not voice_id.strip():
                log_line(log_widget, "X Missing Voice selection."); messagebox.showerror("Error","Choose a voice."); return
            if not os.path.isfile(input_file):
                log_line(log_widget, f"X PowerPoint not found: {input_file}"); messagebox.showerror("Error", f"PowerPoint not found:\n{input_file}"); return
            if not test_ffmpeg_available():
                log_line(log_widget, "X ffmpeg not found."); messagebox.showerror("ffmpeg not found","Place ffmpeg\\ffmpeg.exe next to the EXE or install ffmpeg."); return

            os.makedirs(output_dir, exist_ok=True)
            fixed_dir = output_dir

            log_line(log_widget, f"i Deck: {input_file}")
            log_line(log_widget, f"i Output: {output_dir}")
            log_line(log_widget, "i Loading slides...")

            # Use python-pptx to read notes text
            try:
                prs = Presentation(input_file)
            except Exception as e:
                log_line(log_widget, f"X Failed to open PowerPoint: {e}"); messagebox.showerror("Error", f"Failed to open PowerPoint:\n{e}"); return

            total = len(prs.slides)
            sel = select_slides(total, slide_range_spec or "")
            log_line(log_widget, f"OK Loaded {total} slide(s). Will process: {sel if sel else 'none'}")
            if not sel:
                messagebox.showinfo("No slides selected","Your slide range selected no slides.")
                return

            # Skip PowerPoint operations in audio-only mode
            if audio_only:
                log_line(log_widget, "i Audio-only mode: skipping PowerPoint operations")
                animation_snapshots = {}
            else:
                # Open PowerPoint via COM for animation handling
                log_line(log_widget, "i Opening PowerPoint for animation preservation...")
                try:
                    from win32com.client import Dispatch, GetActiveObject
                    # gencache removed - causes issues in frozen exe
                    
                    # Try to get existing PowerPoint instance first
                    pp_app = None
                    try:
                        pp_app = GetActiveObject("PowerPoint.Application")
                        log_line(log_widget, "  Using existing PowerPoint instance")
                    except:
                        # No existing instance found
                        pass
                    
                    # If no existing instance, create new one
                    if not pp_app:
                        pp_app = Dispatch("PowerPoint.Application")
                        pp_app.Visible = True
                        log_line(log_widget, "  Created new PowerPoint instance")
                    else:
                        pp_app.Visible = True
                    
                    # Check if deck is already open
                    abs_path = os.path.abspath(input_file)
                    pp_pres = None
                    
                    log_line(log_widget, f"  Looking for deck: {os.path.basename(abs_path)}")
                    log_line(log_widget, f"  Currently open presentations: {pp_app.Presentations.Count}")
                    
                    for pres in pp_app.Presentations:
                        try:
                            pres_path = os.path.abspath(pres.FullName)
                            log_line(log_widget, f"  Checking: {os.path.basename(pres_path)}")
                            if pres_path.lower() == abs_path.lower():
                                pp_pres = pres
                                log_line(log_widget, "  Found: Deck already open, reusing")
                                break
                        except Exception as e:
                            log_line(log_widget, f"  Error checking presentation: {e}")
                            continue
                    
                    # If not open, open it
                    if not pp_pres:
                        log_line(log_widget, f"  Opening deck: {os.path.basename(abs_path)}")
                        try:
                            pp_pres = pp_app.Presentations.Open(abs_path, WithWindow=True)
                            log_line(log_widget, "  Deck opened successfully")
                        except Exception as e:
                            log_line(log_widget, f"  Failed to open deck: {e}")
                            raise
                    
                    if not pp_pres:
                        raise RuntimeError("Deck failed to open (pp_pres is None)")
                    
                    log_line(log_widget, "OK PowerPoint ready for animation preservation")
                except Exception as e:
                    log_line(log_widget, f"X Failed to open PowerPoint via COM: {e}")
                    messagebox.showerror("Error", f"Failed to open PowerPoint for animation handling:\n{e}")
                    return

                # BATCH SNAPSHOT: Backup animations for all selected slides upfront
                log_line(log_widget, "i Backing up animations for selected slides...")
                animation_snapshots = {}
                
                for idx in sel:
                    if cancel_event.is_set():
                        break
                    try:
                        slide = pp_pres.Slides(idx)
                        
                        # Clean up orphaned effects first
                        voxanimate.cleanup_orphaned_audio_effects(slide)
                        
                        # Snapshot the animation state
                        snapshot = voxanimate.snapshot_slide_animations(slide)
                        animation_snapshots[idx] = snapshot
                        
                        if snapshot.get("effects"):
                            log_line(log_widget, f"  Slide {idx:02d}: {len(snapshot['effects'])} animations backed up")
                        else:
                            log_line(log_widget, f"  Slide {idx:02d}: No animations")
                            
                    except Exception as e:
                        log_line(log_widget, f"  Slide {idx:02d}: Backup failed - {e}")
                        animation_snapshots[idx] = None
                
                log_line(log_widget, "OK Animation backup complete\n")

            url = f"https://api.elevenlabs.io/v1/text-to-speech/{voice_id}"
            h = {"xi-api-key": api_key, "Content-Type": "application/json"}

            processed = 0

            # Process each slide: TTS generation + audio insertion + animation restoration
            for idx in sel:
                if cancel_event.is_set():
                    log_line(log_widget, "i Run cancelled by user.")
                    break

                # Get notes text from python-pptx
                s = prs.slides[idx-1]
                text = s.notes_slide.notes_text_frame.text if s.notes_slide and s.notes_slide.notes_text_frame else ""
                note = (text or "").strip()

                # Check for "### Read Slide" marker (case-insensitive)
                import re
                read_slide_pattern = re.compile(r'###\s*read\s*slide', re.IGNORECASE)
                if read_slide_pattern.search(note):
                    log_line(log_widget, f"   Detected '### Read Slide' marker - extracting slide text...")
                    try:
                        # Extract text from slide shapes (excluding title)
                        slide_text = extract_slide_text(s)
                        if slide_text:
                            # Replace the marker with extracted text (case-insensitive)
                            note = read_slide_pattern.sub(slide_text, note)
                            log_line(log_widget, f"   Extracted {len(slide_text)} chars from slide")
                        else:
                            log_line(log_widget, f"   Warning: No text found on slide to extract")
                            # Remove the marker so we don't generate audio for it
                            note = read_slide_pattern.sub("", note).strip()
                    except Exception as e:
                        log_line(log_widget, f"   Error extracting slide text: {e}")
                        # Continue with original notes (minus the marker)
                        note = read_slide_pattern.sub("", note).strip()

                if not note:
                    log_line(log_widget, f"- Skipping slide {idx:02d}: No notes found.")
                    processed += 1
                    continue

                log_line(log_widget, f"> Generating slide {idx:02d}...")
                txt_hash = hashlib.sha256(note.encode("utf-8", "ignore")).hexdigest()[:8]
                log_line(log_widget, f"   text#={txt_hash}")
                payload = {"text": note, "output_format": "wav", "voice_settings": {"stability": 0.5, "similarity_boost": 0.7}}

                # Generate TTS audio
                try:
                    sess = VOX_SESSION
                    attempts = 0
                    last_err = None
                    resp = None
                    while attempts < NET_MAX_ATTEMPTS:
                        attempts += 1
                        try:
                            resp = sess.post(url, headers=h, json=payload, timeout=NET_TIMEOUT)
                            if resp.status_code == 200:
                                break
                            if 500 <= resp.status_code < 600 and attempts < NET_MAX_ATTEMPTS:
                                time.sleep(NET_BACKOFF_BASE * (2 ** (attempts-1)))
                                continue
                            break
                        except requests.RequestException as e:
                            last_err = e
                            if attempts < NET_MAX_ATTEMPTS:
                                time.sleep(NET_BACKOFF_BASE * (2 ** (attempts-1)))
                                continue
                            raise

                except requests.RequestException as e:
                    log_line(log_widget, f" X Network error on slide {idx:02d}: {e}")
                    processed += 1
                    continue

                if resp.status_code == 200:
                    name = f"slide{idx:02d}.wav"
                    fixed_path = os.path.join(fixed_dir, name)

                    # Convert audio to proper format
                    with tempfile.NamedTemporaryFile(delete=False, suffix=".wav") as tmp:
                        tmp.write(resp.content); tmp_path = tmp.name
                    wav_md5 = hashlib.md5(resp.content).hexdigest()[:8]
                    log_line(log_widget, f"   wav#={wav_md5}")
                    
                    # Save to manifest
                    try:
                        manifest_path = session_manifest_path
                        try:
                            with open(manifest_path, 'r', encoding='utf-8') as mf:
                                manifest = json.load(mf)
                        except Exception:
                            manifest = []
                        manifest.append({'slide': idx, 'voice_id': voice_id, 'text_sha256': txt_hash, 'wav_md5': wav_md5, 'wav_sha256': _sha256_bytes(resp.content), 'bytes': len(resp.content), 'attempts': attempts, 'http_status': getattr(resp, 'status_code', None)})
                        with open(manifest_path, 'w', encoding='utf-8') as mf:
                            json.dump(manifest, mf, indent=2)
                        try:
                            _msha = _sha256_file(manifest_path)
                            with open(manifest_path + '.sha256', 'w', encoding='utf-8') as csf:
                                csf.write(_msha)
                            logger = logging.getLogger('voxsmith')
                            logger.info(_redact(f'MANIFEST sha256={_msha} file={safe_path(manifest_path)}'))
                        except Exception:
                            pass
                    except Exception:
                        pass
                    
                    run_ffmpeg_quiet(["ffmpeg","-y","-i",tmp_path,"-acodec",TARGET_CODEC,"-ar",TARGET_SAMPLE_RATE,"-ac",TARGET_CHANNELS,fixed_path])
                    try:
                        os.remove(tmp_path)
                    except Exception:
                        pass
                    log_line(log_widget, f" i Converted -> {name}")
                    
                    # Skip attachment if audio_only mode is enabled
                    if audio_only:
                        log_line(log_widget, f"i Audio-only mode: saved to {name}")
                        processed += 1
                        continue
                    
                    # Check if this slide has text animations - if so, skip attachment
                    snapshot = animation_snapshots.get(idx)
                    if snapshot:
                        should_skip, skip_reason = voxanimate.should_skip_audio_attachment(snapshot)
                        if should_skip:
                            log_line(log_widget, f"i Skipping attachment to slide {idx:02d} due to animation backup limitations")
                            log_line(log_widget, f"i {skip_reason}")
                            log_line(log_widget, f"i Audio saved to {name} - attach manually to preserve animations")
                            processed += 1
                            continue

                    
                    # INSERT AUDIO + RESTORE ANIMATIONS
                    log_line(log_widget, f"> Inserting audio into slide {idx:02d}...")
                    
                    try:
                        slide = pp_pres.Slides(idx)
                        
                        # Remove existing VOX audio shapes
                        voxattach._delete_existing_vox_audio(slide)
                        
                        # CRITICAL: Clear animation timeline BEFORE inserting audio
                        # Inserting audio into an animated slide scrambles existing animations
                        log_line(log_widget, f"  Clearing animation timeline before audio insertion...")
                        try:
                            seq = slide.TimeLine.MainSequence
                            while seq.Count > 0:
                                try:
                                    seq.Item(1).Delete()
                                except:
                                    break
                        except Exception:
                            pass
                        
                        # Insert new audio shape WITHOUT auto-creating animation
                        # AddMediaObject (not AddMediaObject2) gives us more control
                        audio_path_abs = os.path.abspath(fixed_path)
                        
                        try:
                            # Try AddMediaObject first (doesn't auto-animate)
                            audio_shape = slide.Shapes.AddMediaObject(audio_path_abs, False, True, 0, 0)
                        except:
                            # Fall back to AddMediaObject2 if AddMediaObject not available
                            audio_shape = slide.Shapes.AddMediaObject2(audio_path_abs, False, True, 0, 0)
                        
                        # Configure audio shape appearance and position
                        try:
                            audio_shape.Width = 32
                            audio_shape.Height = 32
                            W = pp_pres.PageSetup.SlideWidth
                            H = pp_pres.PageSetup.SlideHeight
                            audio_shape.Left = W + 5  # Off-slide to the right
                            audio_shape.Top = H - audio_shape.Height - 5  # Bottom aligned
                            audio_shape.AlternativeText = "VOX_VO"
                            
                            # Disable interactive triggers/click actions
                            try:
                                audio_shape.ActionSettings[1].Action = 0  # ppActionNone
                            except:
                                pass
                        except Exception:
                            pass
                        
                        # RESTORE ANIMATIONS from snapshot
                        snapshot = animation_snapshots.get(idx)
                        if snapshot is not None:
                            # Always call restore - it handles both cases:
                            # 1. If snapshot has effects: restores them with audio at position 1
                            # 2. If snapshot is empty: just adds audio effect
                            log_line(log_widget, f"  Restoring animations...")
                            success = voxanimate.restore_slide_animations(slide, snapshot, audio_shape)
                            if success:
                                effect_count = len(snapshot.get("effects", []))
                                if effect_count > 0:
                                    log_line(log_widget, f"  OK Restored {effect_count} animations")
                                else:
                                    log_line(log_widget, f"  OK Audio inserted (no animations to restore)")
                            else:
                                log_line(log_widget, f"  ! Animation restoration had issues")
                        else:
                            # Snapshot failed, fall back to basic audio setup
                            voxattach._configure_play_settings(audio_shape, hide=True)
                            voxattach._append_media_play_after_previous(slide, audio_shape)
                            log_line(log_widget, f"  ! Snapshot unavailable, basic audio setup used")
                        
                        # Save after each slide
                        pp_pres.Save()
                        log_line(log_widget, f"OK Slide {idx:02d} complete")
                        
                    except Exception as e:
                        log_line(log_widget, f"X Slide {idx:02d} insertion error: {e}")
                        
                else:
                    msg = pretty_api_error(resp)
                    log_line(log_widget, f" X API error slide {idx:02d}: {msg}")

                processed += 1

            if not cancel_event.is_set():
                if audio_only:
                    log_line(log_widget, "* Done. Audio files saved to output folder.")
                    messagebox.showinfo("Complete","Audio generation finished. Files saved to output folder.")
                else:
                    log_line(log_widget, "* Done. Check your output folder.")
                    messagebox.showinfo("Complete","Narration finished. Check your output folder.")
                try:
                    log_line(log_widget, f"OK Session log saved")
                except Exception:
                    pass
            else:
                messagebox.showinfo("Cancelled","Generation was cancelled.")

        except Exception as e:
            log_line(log_widget, f"X Fatal error: {e}")
            traceback.print_exc()
        finally:
            # Save and leave PowerPoint open (don't close) - unless audio_only mode
            if not audio_only:
                try:
                    if pp_pres:
                        pp_pres.Save()
                        log_line(log_widget, "i Deck saved and left open for review")
                except Exception as e:
                    log_line(log_widget, f"! Warning: Failed to save: {e}")
            
            start_button.configure(state="normal")
            cancel_button.configure(state="disabled")
            cancel_event.clear()

    start_button.configure(state="disabled")
    cancel_button.configure(state="normal")
    threading.Thread(target=worker, daemon=True).start()

def main():
    settings = load_settings()

    # Migrate plaintext api_key from settings to OS vault
    if settings.get('api_key'):
        try:
            set_api_key(settings.get('api_key',''))
            settings.pop('api_key', None)
            save_settings(**settings)
        except Exception:
            pass

    # Set CustomTkinter appearance and theme
    ctk.set_appearance_mode("light")
    ctk.set_default_color_theme("blue")
    
    root = ctk.CTk()
    root.configure(fg_color="#e8e8e8")

    try:
        _ok, _msg = verify_self()
        _ver = " - Verified" if _ok else ""
    except Exception:
        _ver = ""
    root.title(f"Voxsmith 2 - {APP_VERSION}{_ver}")

    # --- CustomTkinter Menu Bar ---
    # Create menu bar frame at the top
    menu_frame = ctk.CTkFrame(root, height=40, corner_radius=0, fg_color=("gray90", "gray20"))
    menu_frame.grid(row=0, column=0, sticky="ew", padx=0, pady=0)
    menu_frame.grid_propagate(False)  # Keep fixed height
    
    # We'll add menu buttons after variables are created
    
    # Configure ttk styles for standard tkinter widgets
    style = ttk.Style()
    style.theme_use('clam')  # Use clam theme as base for better styling
    
    # Combobox styling - clean modern look with rounded corners effect
    style.configure('Voxsmith.TCombobox',
                   fieldbackground='#ffffff',
                   background='#ffffff',
                   foreground='#1a1a1a',
                   bordercolor='#d1d5db',
                   arrowcolor='#4a5568',
                   arrowsize=14,
                   padding=(10, 6),
                   relief='flat')
    style.map('Voxsmith.TCombobox',
             fieldbackground=[('readonly', '#ffffff'), ('disabled', '#f3f4f6')],
             foreground=[('readonly', '#1a1a1a'), ('disabled', '#9ca3af')],
             bordercolor=[('focus', '#3b82f6'), ('readonly', '#d1d5db')],
             arrowcolor=[('disabled', '#d1d5db')])
    
    # Configure the dropdown listbox appearance
    root.option_add('*TCombobox*Listbox.background', '#ffffff')
    root.option_add('*TCombobox*Listbox.foreground', '#1a1a1a')
    root.option_add('*TCombobox*Listbox.selectBackground', '#e6f2ff')
    root.option_add('*TCombobox*Listbox.selectForeground', '#1a1a1a')
    root.option_add('*TCombobox*Listbox.font', ('Open Sans', 12))
    
    # Checkbutton styling - much larger, modern checkbox appearance (4x larger)
    style.configure('Voxsmith.TCheckbutton',
                   background='#e8e8e8',
                   foreground='#1a1a1a',
                   font=('Open Sans', 16, 'bold'),
                   borderwidth=0,
                   focuscolor='#e8e8e8',
                   indicatorsize=28,
                   padding=(12, 10),
                   relief='flat')
    style.map('Voxsmith.TCheckbutton',
             background=[('active', '#e8e8e8'), ('pressed', '#e8e8e8')],
             foreground=[('active', '#1a1a1a')],
             indicatorcolor=[('selected', '#3b82f6'), ('!selected', '#ffffff')],
             indicatorrelief=[('pressed', 'sunken')])

    # Main frame with light gray background
    frm = ctk.CTkFrame(root, fg_color="#e8e8e8", corner_radius=0)
    frm.grid(row=1, column=0, sticky="nsew", padx=20, pady=10)
    root.columnconfigure(0, weight=1)
    root.rowconfigure(1, weight=1)
    
    # API var for tracking login state (UI moved to menu bar)
    api_var = tk.StringVar(value=get_api_key())
    acct_state = tk.StringVar(value=("Logged in" if api_var.get().strip() else "Logged out"))

    def do_login():
        # Create custom login dialog matching the UI theme
        login_dialog = ctk.CTkToplevel(root)
        login_dialog.title("Login to Voxsmith")
        login_dialog.configure(fg_color="#e8e8e8")
        login_dialog.geometry("520x300")
        login_dialog.resizable(False, False)
        login_dialog.transient(root)
        login_dialog.grab_set()
        
        # Center the dialog
        login_dialog.update_idletasks()
        x = root.winfo_x() + (root.winfo_width() // 2) - (260)
        y = root.winfo_y() + (root.winfo_height() // 2) - (150)
        login_dialog.geometry(f"+{x}+{y}")
        
        # Content frame
        content = ctk.CTkFrame(login_dialog, fg_color="#e8e8e8", corner_radius=0)
        content.pack(fill="both", expand=True, padx=35, pady=30)
        
        # Title
        ctk.CTkLabel(content, text="Enter your ElevenLabs API Key", 
                    font=("Open Sans", 18, "bold"), text_color="#1a1a1a").pack(pady=(0,8))
        
        # Subtitle with helpful info
        ctk.CTkLabel(content, text="Get your API key from elevenlabs.io/app/settings", 
                    font=("Open Sans", 12), text_color="#555555", 
                    wraplength=420).pack(pady=(0,25))
        
        # API Key entry with placeholder
        key_var = tk.StringVar()
        key_entry = ctk.CTkEntry(content, textvariable=key_var, show="*", 
                                font=("Open Sans", 13), height=44, width=450,
                                placeholder_text="sk-...",
                                fg_color="#ffffff", border_color="#d1d5db", text_color="#1a1a1a")
        key_entry.pack(pady=(0,30))
        key_entry.focus_set()
        
        result = {"submitted": False}
        
        def on_submit():
            result["submitted"] = True
            result["key"] = key_var.get()
            login_dialog.destroy()
        
        def on_cancel():
            login_dialog.destroy()
        
        # Buttons frame
        btn_frame = ctk.CTkFrame(content, fg_color="transparent")
        btn_frame.pack()
        
        ctk.CTkButton(btn_frame, text="Cancel", command=on_cancel, 
                     width=110, height=32, font=("Segoe UI Semibold", 13),
                     fg_color="#e0e0e0", text_color="#333333",
                     hover_color="#d0d0d0", corner_radius=4).pack(side="left", padx=(0,10))
        
        ctk.CTkButton(btn_frame, text="Login", command=on_submit, 
                     width=110, height=32, font=("Segoe UI Semibold", 13, "bold"),
                     corner_radius=4).pack(side="left")
        
        # Bind Enter key
        key_entry.bind("<Return>", lambda e: on_submit())
        login_dialog.bind("<Escape>", lambda e: on_cancel())
        
        # Wait for dialog to close
        root.wait_window(login_dialog)
        
        if result.get("submitted") and result.get("key"):
            key = result["key"].strip()
            if key:
                api_var.set(key)
                set_api_key(key)
                acct_state.set("Logged in")
                # Automatically refresh voices after successful login
                on_refresh_voices()

    def do_logout():
        api_var.set("")
        delete_api_key()
        delete_voice_cache()  # Clear voice cache on logout
        voices_map.clear()  # Clear voices from UI
        voices_combo['values'] = []
        voices_combo.set('')
        voice_id_var.set('')
        acct_state.set("Logged out")
        try:
            remember_var.set(False)
        except Exception:
            pass
        try:
            if os.path.exists(SETTINGS_FILE):
                os.remove(SETTINGS_FILE)
        except Exception:
            pass
        try:
            fb = os.path.join(tempfile.gettempdir(), APP_NAME, 'settings.json')
            if os.path.exists(fb):
                os.remove(fb)
        except Exception:
            pass
    
    def manage_api_key_dialog():
        """Show dialog to manage API key (add/change/remove)"""
        # Create dialog
        dialog = ctk.CTkToplevel(root)
        dialog.title("Manage API Key")
        dialog.configure(fg_color="#e8e8e8")
        dialog.geometry("450x250")
        dialog.resizable(False, False)
        dialog.transient(root)
        dialog.grab_set()
        
        # Center the dialog
        dialog.update_idletasks()
        x = root.winfo_x() + (root.winfo_width() // 2) - (225)
        y = root.winfo_y() + (root.winfo_height() // 2) - (125)
        dialog.geometry(f"+{x}+{y}")
        
        # Content frame
        content = ctk.CTkFrame(dialog, fg_color="#e8e8e8", corner_radius=0)
        content.pack(fill="both", expand=True, padx=30, pady=25)
        
        # Current status
        status_text = "Status: Logged in" if api_var.get().strip() else "Status: Logged out"
        ctk.CTkLabel(content, text=status_text, font=("Open Sans", 14, "bold"), 
                    text_color="#1a1a1a").pack(pady=(0,20))
        
        # Buttons
        if api_var.get().strip():
            # Logged in - show Change and Remove buttons
            ctk.CTkButton(content, text="Change API Key", command=lambda: [dialog.destroy(), do_login()],
                         width=200, height=36, font=("Segoe UI Semibold", 13),
                         corner_radius=4).pack(pady=8)
            ctk.CTkButton(content, text="Remove API Key", command=lambda: [dialog.destroy(), do_logout()],
                         width=200, height=36, font=("Segoe UI Semibold", 13),
                         fg_color="#dc2626", hover_color="#b91c1c",
                         corner_radius=4).pack(pady=8)
        else:
            # Logged out - show Add button
            ctk.CTkButton(content, text="Add API Key", command=lambda: [dialog.destroy(), do_login()],
                         width=200, height=36, font=("Segoe UI Semibold", 13),
                         corner_radius=4).pack(pady=8)
        
        # Close button
        ctk.CTkButton(content, text="Close", command=dialog.destroy,
                     width=200, height=36, font=("Segoe UI Semibold", 13),
                     fg_color="#e0e0e0", text_color="#333333",
                     hover_color="#d0d0d0", corner_radius=4).pack(pady=8)
    
    def show_about_dialog():
        """Show About dialog with version info"""
        # Create dialog
        dialog = ctk.CTkToplevel(root)
        dialog.title("About Voxsmith")
        dialog.configure(fg_color="#e8e8e8")
        dialog.geometry("450x350")
        dialog.resizable(False, False)
        dialog.transient(root)
        dialog.grab_set()
        
        # Center the dialog
        dialog.update_idletasks()
        x = root.winfo_x() + (root.winfo_width() // 2) - (225)
        y = root.winfo_y() + (root.winfo_height() // 2) - (175)
        dialog.geometry(f"+{x}+{y}")
        
        # Content frame
        content = ctk.CTkFrame(dialog, fg_color="#e8e8e8", corner_radius=0)
        content.pack(fill="both", expand=True, padx=30, pady=25)
        
        # Title
        ctk.CTkLabel(content, text=f"Voxsmith {APP_VERSION}", 
                    font=("Open Sans", 22, "bold"), text_color="#1a1a1a").pack(pady=(10,5))
        
        # Main description
        ctk.CTkLabel(content, text="Professional PowerPoint narration with ElevenLabs voices.",
                    font=("Open Sans", 15), text_color="#1a1a1a").pack(pady=5)
        
        # Second line
        ctk.CTkLabel(content, text="This is how you use ElevenLabs with PowerPoint.",
                    font=("Open Sans", 15), text_color="#1a1a1a").pack(pady=5)
        
        # Features
        features = "Preserves animations • Batch processing • Edit mode"
        ctk.CTkLabel(content, text=features,
                    font=("Open Sans", 14), text_color="#1a1a1a").pack(pady=15)
        
        # Built by line
        ctk.CTkLabel(content, text="Built by Don, because he was tired of doing this by hand",
                    font=("Open Sans", 13, "italic"), text_color="#1a1a1a").pack(pady=(10,5))
        
        # Copyright
        import datetime
        year = datetime.datetime.now().year
        ctk.CTkLabel(content, text=f"© {year} Voxsmith",
                    font=("Open Sans", 13), text_color="#1a1a1a").pack(pady=(20,10))
        
        # Close button
        ctk.CTkButton(content, text="Close", command=dialog.destroy,
                     width=120, height=32, font=("Segoe UI Semibold", 13),
                     corner_radius=4).pack(pady=10)

    frm.columnconfigure(1, weight=1)
    frm.columnconfigure(2, weight=1)

    # Deck section header
    ctk.CTkLabel(frm, text="Deck", font=("Open Sans", 16, "bold"), text_color="#1a1a1a").grid(
        row=0, column=0, columnspan=4, sticky="w", pady=(0,10))

    ctk.CTkLabel(frm, text="PowerPoint (.pptx):", font=("Open Sans", 13), text_color="#1a1a1a").grid(row=1, column=0, sticky="w", pady=5, padx=(0,10))
    pptx_var = tk.StringVar(value=settings.get("input_file", DEFAULT_INPUT_FILE))
    pptx_entry = ctk.CTkEntry(frm, textvariable=pptx_var, font=("Open Sans", 13),
                fg_color="#ffffff", border_color="#d1d5db", text_color="#1a1a1a")
    pptx_entry.grid(row=1, column=1, columnspan=2, sticky="ew", pady=5)
    
    def browse_pptx():
        p = filedialog.askopenfilename(
            title="Choose PowerPoint file",
            filetypes=[("PowerPoint files", "*.pptx"), ("All files", "*.*")]
        )
        if p:
            pptx_var.set(p)
            try:
                if remember_var.get():
                    save_settings(
                        voice_id=voice_id_var.get().strip(),
                        voice_name=voices_combo.get().strip(),
                        input_file=pptx_var.get(),
                        output_dir=out_var.get(),
                        preview_text=preview_text_var.get(),
                        fixed_only=True
                    )
            except Exception:
                pass
    
    browse_pptx_btn = ctk.CTkButton(frm, text="Browse...", command=browse_pptx, width=100, height=28, 
                 font=("Segoe UI Semibold", 13), corner_radius=4)
    browse_pptx_btn.grid(row=1, column=3, sticky="e", pady=5, padx=(10,0))
    try:
        browse_pptx_btn.configure(takefocus=0)
    except:
        pass

    ctk.CTkLabel(frm, text="Output folder:", font=("Open Sans", 13), text_color="#1a1a1a").grid(row=2, column=0, sticky="w", pady=5, padx=(0,10))
    out_var = tk.StringVar(value=settings.get("output_dir", DEFAULT_OUTPUT_DIR))
    out_entry = ctk.CTkEntry(frm, textvariable=out_var, font=("Open Sans", 13),
                fg_color="#ffffff", border_color="#d1d5db", text_color="#1a1a1a")
    out_entry.grid(row=2, column=1, columnspan=2, sticky="ew", pady=5)
    
    def browse_out():
        p = filedialog.askdirectory(title="Choose output folder")
        if p:
            out_var.set(p)
    
    browse_out_btn = ctk.CTkButton(frm, text="Browse...", command=browse_out, width=100, height=28,
                 font=("Segoe UI Semibold", 13), corner_radius=4)
    browse_out_btn.grid(row=2, column=3, sticky="e", pady=5, padx=(10,0))
    try:
        browse_out_btn.configure(takefocus=0)
    except:
        pass

    # Audio-only mode checkbox (de-emphasized, out of tab order)
    audio_only_var = tk.BooleanVar(value=bool(settings.get("audio_only", False)))
    audio_only_cb = ctk.CTkCheckBox(frm, text="Generate audio files only (don't attach to slides)", 
                                    variable=audio_only_var, 
                                    font=("Open Sans", 12),
                                    text_color="#6b7280",  # Lighter gray text
                                    fg_color="#17a2b8",
                                    hover_color="#138496",
                                    border_color="#d1d5db",
                                    corner_radius=3)
    audio_only_cb.grid(row=3, column=1, columnspan=2, sticky="w", pady=(2,10))
    # Remove from tab order
    try:
        audio_only_cb.configure(takefocus=0)
    except:
        pass

    # Voice section header
    ctk.CTkLabel(frm, text="Voice", font=("Open Sans", 16, "bold"), text_color="#1a1a1a").grid(
        row=4, column=0, columnspan=4, sticky="w", pady=(20,10))

    ctk.CTkLabel(frm, text="My voices:", font=("Open Sans", 13), text_color="#1a1a1a").grid(row=5, column=0, sticky="w", pady=5, padx=(0,10))
    
    # Using standard ttk.Combobox with custom styling - wrapped in frame for pixel-perfect width
    voice_frame = tk.Frame(frm, bg="#e8e8e8", width=350, height=32)
    voice_frame.grid(row=5, column=1, columnspan=2, sticky="w", pady=5, padx=(0, 5))
    voice_frame.grid_propagate(False)  # Prevent frame from resizing to fit content
    
    voices_combo = ttk.Combobox(voice_frame, values=[], state="readonly",
                                font=("Open Sans", 13), style='Voxsmith.TCombobox')
    voices_combo.place(x=0, y=0, width=350, height=32)  # Fill the frame exactly
    
    voice_id_var = tk.StringVar(value=settings.get("voice_id", ""))

    # Preview Voice button aligned right with Browse buttons
    preview_btn = ctk.CTkButton(frm, text="Preview Voice", font=("Segoe UI Semibold", 13), width=120, height=32,
                               corner_radius=4)
    preview_btn.grid(row=5, column=3, sticky="e", pady=5)

    ctk.CTkLabel(frm, text="Preview text:", font=("Open Sans", 13), text_color="#1a1a1a").grid(row=6, column=0, sticky="w", pady=5, padx=(0,10))
    preview_text_var = tk.StringVar(value=settings.get("preview_text", DEFAULT_PREVIEW_TEXT))
    preview_text_entry = ctk.CTkEntry(frm, textvariable=preview_text_var, font=("Open Sans", 13),
                fg_color="#ffffff", border_color="#d1d5db", text_color="#1a1a1a")
    preview_text_entry.grid(row=6, column=1, columnspan=3, sticky="ew", pady=5)

    # Slide range row
    ctk.CTkLabel(frm, text="Slide range (e.g., 1, 3-6, 10-):", font=("Open Sans", 13), text_color="#1a1a1a").grid(row=7, column=0, sticky="w", pady=(20,5), padx=(0,10))
    slide_range_var = tk.StringVar(value=settings.get("slide_range", DEFAULT_SLIDE_RANGE))
    slide_range_entry = ctk.CTkEntry(frm, textvariable=slide_range_var, font=("Open Sans", 13), width=175,
                fg_color="#ffffff", border_color="#d1d5db", text_color="#1a1a1a")
    slide_range_entry.grid(row=7, column=1, sticky="w", pady=(20,5))

    # Verbose variable (now controlled via Options menu, persistent across sessions)
    verbose_var = tk.BooleanVar(value=bool(settings.get("detailed_logs", False)))

    remember_var = tk.BooleanVar(value=True)
    fixed_only_var = tk.BooleanVar(value=bool(settings.get("fixed_only", DEFAULT_FIXED_ONLY)))

    # Buttons frame - removed Preview Voice and Stop buttons
    btn = ctk.CTkFrame(frm, fg_color="transparent")
    btn.grid(row=8, column=0, columnspan=4, pady=(15,10), sticky="ew")
    
    run_btn = ctk.CTkButton(btn, text="Generate Narration", font=("Segoe UI Semibold", 13), width=140, height=28,
                           corner_radius=4)
    run_btn.grid(row=0, column=0, padx=5)
    
    cancel_btn = ctk.CTkButton(btn, text="Cancel Run", state="disabled", font=("Segoe UI Semibold", 13), width=100, height=28,
                              corner_radius=4)
    cancel_btn.grid(row=0, column=1, padx=5)
    try:
        cancel_btn.configure(takefocus=0)
    except:
        pass
    
    open_out_btn = ctk.CTkButton(btn, text="Open Output", font=("Segoe UI Semibold", 13), width=110, height=28,
                                corner_radius=4)
    open_out_btn.grid(row=0, column=2, padx=5)
    try:
        open_out_btn.configure(takefocus=0)
    except:
        pass
    
    copy_pause_btn = ctk.CTkButton(btn, text="Copy Pause", font=("Segoe UI Semibold", 13), width=100, height=28,
                                  corner_radius=4)
    copy_pause_btn.grid(row=0, column=3, padx=5)
    try:
        copy_pause_btn.configure(takefocus=0)
    except:
        pass

    # Log window with scrollbar
    log_frame = tk.Frame(frm, bg='#ffffff')
    log_frame.grid(row=9, column=0, columnspan=4, pady=(10,0), sticky="nsew")
    
    log = tk.Text(log_frame, height=16, width=100, state="disabled",
                 font=('Consolas', 11), bg='#ffffff', fg='#1a1a1a',
                 relief='flat', borderwidth=0, padx=12, pady=10,
                 wrap=tk.WORD)
    log._verbose_var = verbose_var
    
    scrollbar = tk.Scrollbar(log_frame, command=log.yview, width=16)
    log.configure(yscrollcommand=scrollbar.set)
    
    log.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
    scrollbar.pack(side=tk.RIGHT, fill=tk.Y)

    frm.rowconfigure(9, weight=1)

    voices_map = {}

    def on_choose_voice(event=None):
        choice = voices_combo.get()
        vid = voices_map.get(choice, "")
        if vid:
            voice_id_var.set(vid)
    
    voices_combo.bind("<<ComboboxSelected>>", on_choose_voice)

    def on_refresh_voices():
        k = get_api_key().strip()
        if not k:
            messagebox.showwarning("API Key", "Enter your ElevenLabs API Key first.")
            return
        try:
            items = fetch_voices(k)
        except Exception as e:
            messagebox.showerror("Voices", str(e))
            return
        voices_map.clear(); names = []
        voice_cache_data = []  # For saving to cache
        for nm, vid in items:
            voices_map[nm] = vid; names.append(nm)
            voice_cache_data.append({"name": nm, "voice_id": vid})
        
        # Save to cache for next startup
        save_voice_cache(voice_cache_data)
        
        voices_combo['values'] = names
        saved = load_settings(); sv_vid = saved.get("voice_id", ""); sv_nm = saved.get("voice_name", "")
        sel = None
        if sv_nm and sv_nm in voices_map:
            sel = sv_nm
        elif sv_vid:
            for nm, vid in voices_map.items():
                if vid == sv_vid:
                    sel = nm; break
        if sel:
            voices_combo.set(sel); voice_id_var.set(voices_map[sel])
        elif names:
            voices_combo.set(names[0]); voice_id_var.set(voices_map[names[0]])

    def load_voices_from_cache():
        """Load voices from cache and populate UI immediately."""
        try:
            cache = load_voice_cache()
            if not cache or 'voices' not in cache:
                return False
            
            cached_voices = cache['voices']
            if not cached_voices:
                return False
            
            voices_map.clear()
            names = []
            for v in cached_voices:
                if isinstance(v, dict) and 'name' in v and 'voice_id' in v:
                    nm = v['name']
                    vid = v['voice_id']
                    voices_map[nm] = vid
                    names.append(nm)
            
            if not names:
                return False
                
            voices_combo['values'] = names
            saved = load_settings()
            sv_vid = saved.get("voice_id", "")
            sv_nm = saved.get("voice_name", "")
            sel = None
            if sv_nm and sv_nm in voices_map:
                sel = sv_nm
            elif sv_vid:
                for nm, vid in voices_map.items():
                    if vid == sv_vid:
                        sel = nm
                        break
            if sel:
                voices_combo.set(sel)
                voice_id_var.set(voices_map[sel])
            elif names:
                voices_combo.set(names[0])
                voice_id_var.set(voices_map[names[0]])
            
            return True
        except Exception:
            return False
    
    def background_voice_refresh():
        """Fetch fresh voices in background and update if changed."""
        try:
            k = get_api_key().strip()
            if not k:
                return
            
            # Fetch fresh voices
            fresh_items = fetch_voices(k)
            
            # Get cached voices for comparison
            cache = load_voice_cache()
            cached_voices = cache.get('voices', []) if cache else []
            
            # Check if voices changed
            if voices_changed(cached_voices, fresh_items):
                # Update cache
                voice_cache_data = [{"name": nm, "voice_id": vid} for nm, vid in fresh_items]
                save_voice_cache(voice_cache_data)
                
                # Update UI on main thread
                def update_ui():
                    voices_map.clear()
                    names = []
                    for nm, vid in fresh_items:
                        voices_map[nm] = vid
                        names.append(nm)
                    
                    voices_combo['values'] = names
                    # Try to preserve current selection if it still exists
                    current_sel = voices_combo.get()
                    if current_sel and current_sel in voices_map:
                        voices_combo.set(current_sel)
                        voice_id_var.set(voices_map[current_sel])
                    elif names:
                        # Fall back to saved or first voice
                        saved = load_settings()
                        sv_vid = saved.get("voice_id", "")
                        sv_nm = saved.get("voice_name", "")
                        sel = None
                        if sv_nm and sv_nm in voices_map:
                            sel = sv_nm
                        elif sv_vid:
                            for nm, vid in voices_map.items():
                                if vid == sv_vid:
                                    sel = nm
                                    break
                        if sel:
                            voices_combo.set(sel)
                            voice_id_var.set(voices_map[sel])
                        else:
                            voices_combo.set(names[0])
                            voice_id_var.set(voices_map[names[0]])
                
                root.after(0, update_ui)
            else:
                # Voices unchanged, just update timestamp
                voice_cache_data = [{"name": nm, "voice_id": vid} for nm, vid in fresh_items]
                save_voice_cache(voice_cache_data)
        except Exception:
            # Silently fail - user already has cached voices
            pass

    preview_player = PreviewPlayer(log_widget=log, preview_btn=preview_btn, stop_btn=None,
                                   get_preview_text=lambda: preview_text_var.get())

    cancel_event = threading.Event()

    def on_run():
        vn = voices_combo.get().strip(); vid = voice_id_var.get().strip()
        if remember_var.get():
            save_settings(voice_id=vid, voice_name=vn, input_file=pptx_var.get(), output_dir=out_var.get(),
                          preview_text=preview_text_var.get(), fixed_only=True, audio_only=audio_only_var.get())
        generate_narration(
            api_key=get_api_key(),
            voice_id=vid,
            input_file=pptx_var.get(),
            output_dir=out_var.get(),
            fixed_only=True,
            slide_range_spec=slide_range_var.get(),
            cancel_event=cancel_event,
            log_widget=log,
            start_button=run_btn,
            cancel_button=cancel_btn,
            audio_only=audio_only_var.get()
        )

    def on_cancel():
        cancel_event.set(); cancel_btn.configure(state="disabled"); log._verbose_var = verbose_var; log_line(log, "i Cancelling after current slide...")

    def on_preview():
        if remember_var.get():
            save_settings(preview_text=preview_text_var.get())
        preview_player.preview(api_key=get_api_key(), voice_id=voice_id_var.get())

    def on_open_output():
        p = out_var.get().strip()
        if p and os.path.isdir(p):
            open_folder(p)
        else:
            messagebox.showinfo("Open Output", "Choose a valid output folder first.")

    run_btn.configure(command=on_run)
    cancel_btn.configure(command=on_cancel)
    preview_btn.configure(command=on_preview)
    
    def on_copy_pause():
        try:
            root.clipboard_clear()
            root.clipboard_append('<break time="1s" />')
            root.update_idletasks()
            messagebox.showinfo("Copied!", "Pause code copied to clipboard!")
        except Exception as e:
            messagebox.showerror("Error", f"Failed to copy to clipboard:\n{e}")

    open_out_btn.configure(command=on_open_output)
    copy_pause_btn.configure(command=on_copy_pause)
    
    # Add tooltips to buttons
    ToolTip(preview_btn, "Test the selected voice with preview text")
    ToolTip(run_btn, "Generate audio for selected slides and insert into deck")
    ToolTip(cancel_btn, "Stop the current generation process")
    ToolTip(open_out_btn, "Open the output folder in File Explorer")
    ToolTip(copy_pause_btn, "Copy 1-second break code to clipboard for PowerPoint notes")
    
    # Bind Enter key to slide range entry to trigger Generate
    slide_range_entry.bind("<Return>", lambda e: run_btn.invoke() if slide_range_var.get().strip() else None)
    
    # Configure CustomTkinter menu bar
    def show_options_popup(event=None):
        """Show options popup menu."""
        popup = tk.Menu(root, tearoff=0, font=("Open Sans", 13))
        popup.add_command(label="Manage API Key...", command=manage_api_key_dialog, font=("Open Sans", 13))
        
        def toggle_detailed_logs():
            """Toggle detailed logs and save to settings."""
            save_settings(detailed_logs=verbose_var.get())
        
        popup.add_checkbutton(label="Detailed Logs", variable=verbose_var, 
                             command=toggle_detailed_logs, font=("Open Sans", 13))
        
        try:
            # Position popup below the Options button
            popup.tk_popup(event.x_root, event.y_root + 10)
        finally:
            popup.grab_release()
    
    # Options button
    options_btn = ctk.CTkButton(
        menu_frame,
        text="Options",
        width=80,
        height=30,
        fg_color="transparent",
        text_color=("gray10", "gray90"),
        hover_color=("gray80", "gray30"),
        font=("Open Sans", 13),
        command=lambda: show_options_popup(type('Event', (), {'x_root': options_btn.winfo_rootx(), 'y_root': options_btn.winfo_rooty() + options_btn.winfo_height()})())
    )
    options_btn.pack(side="left", padx=10, pady=5)
    
    # About button
    about_btn = ctk.CTkButton(
        menu_frame,
        text="About",
        width=70,
        height=30,
        fg_color="transparent",
        text_color=("gray10", "gray90"),
        hover_color=("gray80", "gray30"),
        font=("Open Sans", 13),
        command=show_about_dialog
    )
    about_btn.pack(side="left", padx=5, pady=5)
    
    # Account status label on the right
    account_status_label = ctk.CTkLabel(
        menu_frame,
        text=acct_state.get(),
        text_color=("gray30", "gray70"),
        font=("Open Sans", 13)
    )
    account_status_label.pack(side="right", padx=15)
    
    # Trace account state changes to update menu bar label
    def on_account_state_change(*args):
        """Update the account status label in menu bar."""
        account_status_label.configure(text=acct_state.get())
    acct_state.trace_add("write", on_account_state_change)
    
    # Set explicit tab order for keyboard navigation
    # Order: PPTX -> Output folder -> My voices -> Preview text -> Preview Voice -> Slide Range -> Generate Narration -> (loop to PPTX)
    try:
        # Set tab order using lift (raises widget in stacking order for tab traversal)
        # Later widgets lifted = earlier in tab order (reverse order)
        run_btn.lift()
        slide_range_entry.lift()
        preview_btn.lift()
        preview_text_entry.lift()
        voices_combo.lift()
        out_entry.lift()
        pptx_entry.lift()
    except Exception:
        pass  # Tab order is nice-to-have, don't fail if it doesn't work

    saved = load_settings()
    if saved.get("api_key", "").strip():
        try:
            items = fetch_voices(saved.get("api_key").strip()); names = []
            voices_map.clear()
            for nm, vid in items:
                voices_map[nm] = vid; names.append(nm)
            voices_combo['values'] = names
            sv_vid = saved.get("voice_id", ""); sv_nm = saved.get("voice_name", ""); sel = None
            if sv_nm and sv_nm in voices_map:
                sel = sv_nm
            elif sv_vid:
                for nm, vid in voices_map.items():
                    if vid == sv_vid:
                        sel = nm; break
            if sel:
                voices_combo.set(sel); voice_id_var.set(voices_map[sel])
        except Exception:
            pass

    def on_close():
        vn = voices_combo.get().strip(); vid = voice_id_var.get().strip()
        save_settings(voice_id=vid, voice_name=vn, input_file=pptx_var.get(), output_dir=out_var.get(),
                      preview_text=preview_text_var.get(), fixed_only=True)
        try:
            _release_single_instance(_SINGLE_LOCK)
        except Exception:
            pass
        root.destroy()

    root.protocol("WM_DELETE_WINDOW", on_close)
    root.minsize(900, 600)
    
    # Load voices asynchronously in background after window appears
    def async_load_voices_on_startup():
        """Load voices with smart caching: instant cache + background refresh."""
        try:
            if not _is_logged_out():
                k = get_api_key()
                if k and k.strip():
                    # Step 1: Try to load from cache immediately (instant)
                    cache_loaded = load_voices_from_cache()
                    
                    # Step 2: ALWAYS refresh in background (even if cache valid)
                    def start_background_refresh():
                        threading.Thread(
                            target=background_voice_refresh,
                            daemon=True
                        ).start()
                    
                    # If cache loaded, start background refresh after short delay
                    # If no cache, do foreground refresh
                    if cache_loaded:
                        root.after(150, start_background_refresh)
                    else:
                        # No cache, do normal refresh
                        root.after(100, lambda: threading.Thread(
                            target=lambda: root.after(0, on_refresh_voices),
                            daemon=True
                        ).start())
        except Exception:
            pass
    
    # Schedule async voice loading after window appears
    root.after(50, async_load_voices_on_startup)

    root.mainloop()

if __name__ == "__main__":
    try:
        try:
            already_open, _SINGLE_LOCK = _check_single_instance()
        except Exception as e:
            already_open, _SINGLE_LOCK = False, None
            print(f"Single-instance check failed: {e}")

        if already_open:
            try:
                r = tk.Tk(); r.withdraw()
                messagebox.showinfo("Already running", "The app is already open.")
                r.destroy()
            except Exception:
                pass
            sys.exit(0)

        main()

    except Exception:
        try:
            log_dir = os.path.join(os.getenv("APPDATA") or tempfile.gettempdir(), APP_NAME)
            os.makedirs(log_dir, exist_ok=True)
            log_path = os.path.join(log_dir, "latest.log")
            with open(log_path, "w", encoding="utf-8") as f:
                f.write(traceback.format_exc())
            try:
                r = tk.Tk(); r.withdraw()
                messagebox.showerror(f"{APP_NAME} - Error", f"An error occurred while starting the app.\n\nDetails were written to:\n{log_path}")
                r.destroy()
            except Exception:
                pass
        finally:
            try:
                _release_single_instance(_SINGLE_LOCK)
            except Exception:
                pass

def _handle_auth_log(resp):
    try:
        s = getattr(resp, "status_code", None)
        if s in (401, 403):
            logging.getLogger("voxsmith").info(_redact("AUTH issue: 401/403 from ElevenLabs. Check API key in Credential Manager."))
    except Exception:
        pass
